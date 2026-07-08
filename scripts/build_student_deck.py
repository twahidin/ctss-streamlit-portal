"""Generate the student collaboration guide PPTX.

v2 — per-feature workflow: each member owns one .py module; the Leader owns
the glue (app.py). Includes a slide on Streamlit Playground for quick widget
tests.
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

OUTPUT = Path.home() / "Desktop" / "CTSS_Streamlit_Group_Guide.pptx"

# Theme
BG = RGBColor(0xFB, 0xFB, 0xFD)
INK = RGBColor(0x14, 0x18, 0x24)
MUTED = RGBColor(0x55, 0x5B, 0x6B)
ACCENT = RGBColor(0x4A, 0x52, 0xE0)
ACCENT_SOFT = RGBColor(0xEC, 0xEE, 0xFD)
CODE_BG = RGBColor(0x10, 0x14, 0x1F)
CODE_TXT = RGBColor(0xE8, 0xEA, 0xF2)
WARN = RGBColor(0xC0, 0x4A, 0x1F)
SUCCESS_BG = RGBColor(0xE6, 0xF7, 0xEC)
SUCCESS_TXT = RGBColor(0x1A, 0x6B, 0x35)

SLIDE_W, SLIDE_H = Inches(13.333), Inches(7.5)


def new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def add_blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def fill_bg(slide):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG
    bg.shadow.inherit = False
    return bg


def add_text(slide, x, y, w, h, text, *, size=18, bold=False, color=INK,
             font="Calibri", align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tx = slide.shapes.add_textbox(x, y, w, h)
    tf = tx.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return tx


def add_title(slide, title, subtitle=None, page=None, total=None):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.55), Inches(0.18), Inches(0.65))
    bar.line.fill.background()
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.shadow.inherit = False
    add_text(slide, Inches(0.85), Inches(0.45), Inches(11), Inches(0.85), title, size=32, bold=True)
    if subtitle:
        add_text(slide, Inches(0.85), Inches(1.20), Inches(11), Inches(0.5), subtitle, size=15, color=MUTED)
    if page is not None and total is not None:
        add_text(slide, Inches(11.8), Inches(7.05), Inches(1.3), Inches(0.3),
                 f"{page} / {total}", size=10, color=MUTED, align=PP_ALIGN.RIGHT)
    add_text(slide, Inches(0.5), Inches(7.05), Inches(8), Inches(0.3),
             "CTSS — Streamlit Class Project", size=10, color=MUTED)


def add_bullets(slide, x, y, w, h, items, *, size=18, color=INK):
    tx = slide.shapes.add_textbox(x, y, w, h)
    tf = tx.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(8)
        r1 = p.add_run()
        r1.text = "•  "
        r1.font.size = Pt(size)
        r1.font.bold = True
        r1.font.color.rgb = ACCENT
        r1.font.name = "Calibri"
        r2 = p.add_run()
        r2.text = item
        r2.font.size = Pt(size)
        r2.font.color.rgb = color
        r2.font.name = "Calibri"
    return tx


def add_code(slide, x, y, w, h, code, *, size=14):
    bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = CODE_BG
    bg.shadow.inherit = False
    tx = slide.shapes.add_textbox(x + Inches(0.18), y + Inches(0.10), w - Inches(0.36), h - Inches(0.20))
    tf = tx.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0)
    tf.margin_top = Inches(0)
    lines = code.splitlines() or [""]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(0)
        r = p.add_run()
        r.text = line if line else " "
        r.font.name = "Consolas"
        r.font.size = Pt(size)
        r.font.color.rgb = CODE_TXT
    return bg


def add_callout(slide, x, y, w, h, text, *, color=ACCENT_SOFT, ink=INK, size=14, bold=False):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    box.line.fill.background()
    box.fill.solid()
    box.fill.fore_color.rgb = color
    box.shadow.inherit = False
    tx = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.1), w - Inches(0.4), h - Inches(0.2))
    tf = tx.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0)
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.name = "Calibri"
    r.font.size = Pt(size)
    r.font.color.rgb = ink
    r.font.bold = bold


def add_table(slide, x, y, w, h, headers, rows):
    cols = len(headers)
    n_rows = len(rows) + 1
    tbl_shape = slide.shapes.add_table(n_rows, cols, x, y, w, h)
    table = tbl_shape.table
    for i, htxt in enumerate(headers):
        cell = table.cell(0, i)
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT
        cell.text_frame.text = ""
        p = cell.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = htxt
        r.font.bold = True
        r.font.size = Pt(14)
        r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        r.font.name = "Calibri"
    for ri, row in enumerate(rows, start=1):
        for ci, val in enumerate(row):
            cell = table.cell(ri, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF) if ri % 2 else ACCENT_SOFT
            cell.text_frame.text = ""
            p = cell.text_frame.paragraphs[0]
            r = p.add_run()
            r.text = val
            r.font.size = Pt(13)
            r.font.color.rgb = INK
            r.font.name = "Calibri"


# ----------------- Slides -----------------

def slide_title(prs, total):
    s = add_blank(prs)
    fill_bg(s)
    block = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.4), SLIDE_W, Inches(2.7))
    block.line.fill.background()
    block.fill.solid()
    block.fill.fore_color.rgb = ACCENT
    block.shadow.inherit = False
    add_text(s, Inches(0.5), Inches(2.55), Inches(12), Inches(0.5),
             "CLEMENTI TOWN SECONDARY SCHOOL", size=14, bold=True,
             color=RGBColor(0xC9, 0xCD, 0xF8))
    add_text(s, Inches(0.5), Inches(3.0), Inches(12), Inches(1.2),
             "Building Your Streamlit App\nas a Team", size=44, bold=True,
             color=RGBColor(0xFF, 0xFF, 0xFF))
    add_text(s, Inches(0.5), Inches(4.5), Inches(12), Inches(0.7),
             "Each member owns one feature file. The Leader glues them together.", size=18,
             color=RGBColor(0xE3, 0xE6, 0xFD))
    add_text(s, Inches(0.5), Inches(6.7), Inches(12), Inches(0.4),
             f"A {total}-step playbook for student groups", size=13, color=MUTED)


def slide_workflow(prs, page, total):
    s = add_blank(prs)
    fill_bg(s)
    add_title(s, "The workflow at a glance",
              "One file per person + one Leader to glue it all together.", page, total)
    # Diagram-ish
    add_code(s, Inches(0.85), Inches(1.95), Inches(11.6), Inches(4.0), """  Alice's Colab           Bob's Colab            Carol's Colab           Dave's Colab
  ─────────────           ───────────            ─────────────           ────────────
   weather_loader.py       chart.py              form.py                 styles.py

                                       │
                                       ▼
                                Leader's Colab
                                ──────────────
                                  app.py   ←  glue: imports each feature

                                       │
                                       ▼
                              Upload Portal  ──►  GitHub  ──►  Railway  ──►  Live URL""", size=14)
    add_callout(
        s, Inches(0.85), Inches(6.2), Inches(11.6), Inches(0.95),
        "Each member writes ONE Python file.  The Leader writes app.py and uploads everything together.",
        color=ACCENT_SOFT, size=14,
    )


def slide_roles(prs, page, total):
    s = add_blank(prs)
    fill_bg(s)
    add_title(s, "Pick your roles", "Each Coder owns ONE file. That's their name on it.", page, total)
    add_table(
        s, Inches(0.85), Inches(1.85), Inches(11.6), Inches(3.2),
        ["Role", "How many", "What they own"],
        [
            ["🎯  Leader", "1 (fixed)", "app.py — imports each feature, builds the UI flow"],
            ["💻  Coder A", "1", "One feature file (e.g. weather_loader.py)"],
            ["💻  Coder B", "1", "One feature file (e.g. chart.py)"],
            ["💻  Coder C", "1", "One feature file (e.g. form.py)"],
            ["🧪  Tester", "rotates", "Opens the live URL, files bugs in the bug-list cell"],
        ],
    )
    add_callout(
        s, Inches(0.85), Inches(5.4), Inches(11.6), Inches(1.4),
        "If your group has 3 members: Leader can also be a Coder. "
        "If you have 5 members: split a feature in two. The point is one owner per file — never two.",
        color=ACCENT_SOFT, size=15,
    )


def slide_setup(prs, page, total):
    s = add_blank(prs)
    fill_bg(s)
    add_title(s, "Set up your notebooks",
              "Each person makes their OWN Colab. Then share each link with the Leader.", page, total)
    add_bullets(s, Inches(0.85), Inches(1.85), Inches(11.5), Inches(2.7), [
        "Each member: colab.research.google.com  →  New Notebook.",
        "Name it after your feature — e.g. 'Alice — weather_loader'.",
        "Top-right Share  →  'Anyone with the link can VIEW'  →  copy URL  →  drop in chat.",
        "The Leader makes ONE more notebook called 'Group N — app.py'.",
    ], size=17)
    add_callout(
        s, Inches(0.85), Inches(5.3), Inches(11.6), Inches(1.6),
        "Why one notebook per person?\n"
        "   ☐  No real-time editing chaos.\n"
        "   ☐  Each student practises Python modules — designing functions to be IMPORTED.\n"
        "   ☐  When Alice changes her function signature, the Leader notices fast.",
        color=ACCENT_SOFT, size=14,
    )


def slide_plan(prs, page, total):
    s = add_blank(prs)
    fill_bg(s)
    add_title(s, "Plan together",
              "Decide who owns what BEFORE anyone codes.", page, total)
    add_bullets(s, Inches(0.85), Inches(1.95), Inches(11.5), Inches(2.5), [
        "Open a shared Google Doc (or use the Leader's notebook as a planning doc).",
        "Write what your app does in ONE sentence.",
        "List 3–5 features. Put each owner's INITIALS next to their feature.",
        "Out of scope: 1–2 things you are NOT building. Protects you from feature creep.",
    ], size=17)
    add_code(s, Inches(0.85), Inches(4.7), Inches(11.6), Inches(2.2), """## Weather mood board for Clementi

Features:
- AB  load_weather()        →  weather_loader.py
- SK  show_today_emoji()    →  mood.py
- KC  plot_last_7_days()    →  chart.py
- JT  add_personal_note()   →  notes.py

Out of scope: login, mobile app""")


def slide_naming(prs, page, total):
    s = add_blank(prs)
    fill_bg(s)
    add_title(s, "Naming your feature file",
              "Name it after the FEATURE, not the person.", page, total)
    add_bullets(s, Inches(0.85), Inches(1.85), Inches(11.5), Inches(1.6), [
        "Lowercase letters, numbers, underscores. End with .py.",
        "Pattern:  what_it_does.py   ✅",
        "Avoid:  alice.py, file1.py, my_helper.py   ❌",
    ], size=18)
    add_code(s, Inches(0.85), Inches(4.0), Inches(11.6), Inches(2.6), """  Bad name              Better name             Why
  ───────────────────   ─────────────────────   ─────────────────────────────
  alice.py              weather_loader.py        Names a feature, not a person
  helper.py             chart.py                 Says what it does
  file1.py              form.py                  Reads naturally in app.py
  weatherLoader.py      weather_loader.py        Python style: snake_case""", size=14)


def slide_glue(prs, page, total):
    s = add_blank(prs)
    fill_bg(s)
    add_title(s, "How app.py glues your features",
              "The Leader's job: import each feature and call it.", page, total)
    add_text(s, Inches(0.85), Inches(1.85), Inches(11.5), Inches(0.4),
             "Each feature file (Coder writes this):", size=14, color=MUTED)
    add_code(s, Inches(0.85), Inches(2.25), Inches(11.6), Inches(1.5), """# weather_loader.py   |   owner: AB
import requests

def load_weather(city):
    return requests.get(f\"https://api.example.com/{city}\").json()""")
    add_text(s, Inches(0.85), Inches(3.95), Inches(11.5), Inches(0.4),
             "app.py imports each feature and uses it (Leader writes this):", size=14, color=MUTED)
    add_code(s, Inches(0.85), Inches(4.35), Inches(11.6), Inches(2.4), """# app.py   |   Leader's glue
import streamlit as st
from weather_loader import load_weather
from mood import show_today_emoji
from chart import plot_last_7_days

st.title(\"Clementi Weather Mood\")
data = load_weather(\"Clementi\")
show_today_emoji(data)
plot_last_7_days(data)""")


def slide_code_your_feature(prs, page, total):
    s = add_blank(prs)
    fill_bg(s)
    add_title(s, "Code your feature in your own notebook",
              "You're a writer of one Python module. That's it.", page, total)
    add_bullets(s, Inches(0.85), Inches(1.85), Inches(11.5), Inches(2.0), [
        "Top of your notebook — one cell saves your code as a real .py file:",
    ], size=16)
    add_code(s, Inches(0.85), Inches(2.65), Inches(11.6), Inches(2.5), """%%writefile chart.py
import streamlit as st
import matplotlib.pyplot as plt

def plot_last_7_days(data):
    fig, ax = plt.subplots()
    ax.plot(data[\"days\"], data[\"temps\"])
    st.pyplot(fig)""")
    add_callout(
        s, Inches(0.85), Inches(5.4), Inches(11.6), Inches(1.5),
        "Rules for your file:\n"
        "   ☐  Wrap your code in FUNCTIONS — so app.py can call them.\n"
        "   ☐  Don't put st.title(...) at the top level — that belongs in app.py.\n"
        "   ☐  Test your function in a scratch cell before sharing your link.",
        color=ACCENT_SOFT, size=14,
    )


def slide_buglist(prs, page, total):
    s = add_blank(prs)
    fill_bg(s)
    add_title(s, "Shared bug-list cell",
              "Where? In the LEADER's app.py notebook, last text cell.", page, total)
    add_code(s, Inches(0.85), Inches(1.85), Inches(11.6), Inches(3.4), """# 🐛 Bug list + notes
Format:  [STATUS] file — description — reporter → owner

- [OPEN]   chart.py    wrong colour on dark mode    JT → SK
- [DOING]  form.py     submit doesn't clear input   KC → KC
- [DONE]   weather_loader.py  crash on empty city   fixed by AB ✅

## Tester rotation
Lesson 1 → JT     Lesson 2 → KC     Lesson 3 → SK""")
    add_callout(
        s, Inches(0.85), Inches(5.5), Inches(11.6), Inches(1.4),
        "Rules:  No bugs in chat — they go HERE.  Always include the FILENAME so the right person picks it up.",
        color=ACCENT_SOFT, size=14,
    )


def slide_streamlit_playground(prs, page, total):
    s = add_blank(prs)
    fill_bg(s)
    add_title(s, "Streamlit: playground, docs, tutorials",
              "Three official resources that will save you hours.", page, total)
    add_bullets(s, Inches(0.85), Inches(1.85), Inches(11.5), Inches(3.4), [
        "PLAYGROUND — streamlit.io/playground   ▸  Paste a snippet, hit Run, see the widget live. "
        "Best for: testing single widgets (st.button, st.slider, st.selectbox).",
        "API REFERENCE — docs.streamlit.io/develop/api-reference   ▸  Every widget, with examples. "
        "Use it like a dictionary: 'how do I make a date picker?' → search st.date_input.",
        "TUTORIALS — docs.streamlit.io/develop/tutorials   ▸  Walkthroughs of real apps. "
        "Read one before you start your feature — there's probably a pattern that matches.",
    ], size=14)
    add_code(s, Inches(0.85), Inches(5.4), Inches(11.6), Inches(1.5), """# Try this in the Playground:
import streamlit as st
n = st.slider(\"How many balloons?\", 1, 100, 10)
if st.button(\"Drop them\"):
    for _ in range(n):
        st.balloons()""", size=13)


def slide_pure_python(prs, page, total):
    s = add_blank(prs)
    fill_bg(s)
    add_title(s, "Test pure Python in Colab",
              "If your function doesn't call st.something — test it here first.", page, total)
    add_code(s, Inches(0.85), Inches(1.95), Inches(11.6), Inches(3.0), """# In your scratch cell, test the pure logic:
def average(numbers):
    return sum(numbers) / len(numbers)

print(average([3, 5, 8]))    # 5.33
print(average([10, 20]))     # 15.0""")
    add_callout(
        s, Inches(0.85), Inches(5.2), Inches(11.6), Inches(1.4),
        "Rule of thumb:  pure Python  →  test in Colab.  "
        "Streamlit widgets  →  test in the Playground.  "
        "Full app  →  upload through the portal.",
        color=ACCENT_SOFT, size=14,
    )


def slide_requirements(prs, page, total):
    s = add_blank(prs)
    fill_bg(s)
    add_title(s, "Pick your libraries (no requirements.txt!)",
              "The portal now has checkboxes. Tick what your code imports.", page, total)
    add_bullets(s, Inches(0.85), Inches(1.95), Inches(11.5), Inches(3.0), [
        "Open every .py file in your group. List every 'import' line.",
        "On the portal, tick the boxes that match (streamlit is always ticked).",
        "Anything not in the list? Add it in the 'Other libraries' textarea, one per line.",
        "Built-ins to SKIP (already in Python): os, sys, re, json, math, random, datetime, time.",
    ], size=16)
    add_callout(
        s, Inches(0.85), Inches(5.3), Inches(11.6), Inches(1.4),
        "Common ones available as checkboxes: pandas, numpy, matplotlib, plotly, requests, "
        "anthropic, openai, scikit-learn, pillow, streamlit-extras.",
        color=ACCENT_SOFT, size=14,
    )


def slide_access(prs, page, total):
    s = add_blank(prs)
    fill_bg(s)
    add_title(s, "How to access the portal & use the system",
              "One web page. Group code in, live app out.", page, total)
    add_bullets(s, Inches(0.85), Inches(1.85), Inches(11.5), Inches(2.2), [
        "Open a browser  →  go to  https://your-portal.up.railway.app",
        "Type your group's CODE (the 'team password' on the next slides)  →  Continue.",
        "Now you see the upload form: app.py, Feature files 1–4, and library checkboxes.",
        "Fill it in (see 'Leader ships everything')  →  click Deploy  →  wait ~1–3 min.",
        "Visit your group's live URL (e.g. https://group3.up.railway.app) — that's your app.",
    ], size=16)
    add_code(s, Inches(0.85), Inches(4.85), Inches(11.6), Inches(1.25), """  https://your-portal.up.railway.app   →   [ enter group code ]   →   upload   →   Deploy
                                                                          │
                                              your live app   ◄───────────┘   https://grpN.up.railway.app""", size=13)
    add_callout(
        s, Inches(0.85), Inches(6.35), Inches(11.6), Inches(0.8),
        "Keep your group code secret — anyone with it can overwrite your app. Tell your teacher if it leaks.",
        color=ACCENT_SOFT, size=14,
    )


def _passwords_slide(prs, page, total, title, rows):
    s = add_blank(prs)
    fill_bg(s)
    add_title(s, title,
              "Portal:  https://your-portal.up.railway.app   —   keep your code secret.", page, total)
    add_table(
        s, Inches(0.85), Inches(1.95), Inches(11.6), Inches(0.5 + 0.5 * len(rows)),
        ["Group", "Team password (group code)", "Your live app URL"],
        rows,
    )
    add_callout(
        s, Inches(0.85), Inches(6.5), Inches(11.6), Inches(0.75),
        "Only the Leader needs to type this into the portal. The live URL is what you share with everyone.",
        color=ACCENT_SOFT, size=14,
    )


# NOTE: The codes below are PLACEHOLDERS. Generate real ones per deployment with
#   python -c "import secrets; print(secrets.token_urlsafe(8))"
# and keep them out of version control (they are the live "team passwords").
def slide_passwords_3_8(prs, page, total):
    _passwords_slide(prs, page, total, "Team passwords — Groups 3 to 8", [
        ["Group 3", "<group-3-code>", "https://group3.up.railway.app"],
        ["Group 4", "<group-4-code>", "https://group4.up.railway.app"],
        ["Group 5", "<group-5-code>", "https://group5.up.railway.app"],
        ["Group 6", "<group-6-code>", "https://group6.up.railway.app"],
        ["Group 7", "<group-7-code>", "https://group7.up.railway.app"],
        ["Group 8", "<group-8-code>", "https://group8.up.railway.app"],
    ])


def slide_passwords_9_11(prs, page, total):
    _passwords_slide(prs, page, total, "Team passwords — Groups 9 to 11", [
        ["Group 9",  "<group-9-code>",  "https://group9.up.railway.app"],
        ["Group 10", "<group-10-code>", "https://group10.up.railway.app"],
        ["Group 11", "<group-11-code>", "https://group11.up.railway.app"],
    ])


def slide_passwords_12_16(prs, page, total):
    _passwords_slide(prs, page, total, "Team passwords — Groups 12 to 16", [
        ["Group 12", "<group-12-code>", "https://group12.up.railway.app"],
        ["Group 13", "<group-13-code>", "https://group13.up.railway.app"],
        ["Group 14", "<group-14-code>", "https://group14.up.railway.app"],
        ["Group 15", "<group-15-code>", "https://group15.up.railway.app"],
        ["Group 16", "<group-16-code>", "https://group16.up.railway.app"],
    ])


def slide_ship(prs, page, total):
    s = add_blank(prs)
    fill_bg(s)
    add_title(s, "Leader ships everything",
              "One trip to the portal — uploads or pastes every link.", page, total)
    add_bullets(s, Inches(0.85), Inches(1.85), Inches(11.5), Inches(3.4), [
        "Open the upload portal:  https://your-portal.up.railway.app",
        "Enter your group code.",
        "app.py: pick a file OR paste your Leader notebook link.",
        "Open Feature file 1 → name it (e.g. weather_loader.py) → paste Alice's notebook link.",
        "Repeat for Feature 2, 3, 4 (one per Coder who has a file ready today).",
        "Tick the libraries used, then click Deploy.",
    ], size=16)
    add_callout(
        s, Inches(0.85), Inches(5.7), Inches(11.6), Inches(1.2),
        "Each notebook needs sharing set to 'Anyone with the link – Viewer'. "
        "The portal extracts the cell that starts with %%writefile <your-filename>.",
        color=ACCENT_SOFT, size=14,
    )


def slide_iterate(prs, page, total):
    s = add_blank(prs)
    fill_bg(s)
    add_title(s, "Iterate", "Real software is never 'done' the first time.", page, total)
    add_bullets(s, Inches(0.85), Inches(1.85), Inches(11.5), Inches(3.0), [
        "Goal: ship 3+ versions in this lesson.",
        "Cycle: Tester finds a bug → Coder fixes IN THEIR NOTEBOOK → Leader re-uploads.",
        "Don't all upload at once — Leader does ONE trip with everyone's latest links.",
        "Done is NOT 'the code runs'. Done is 'Tester signed off in the bug list'.",
    ], size=17)
    add_callout(
        s, Inches(0.85), Inches(5.0), Inches(11.6), Inches(1.7),
        "End-of-lesson checklist:\n"
        "   ☐  Every member's notebook has a working %%writefile <feature>.py cell\n"
        "   ☐  All [OPEN] bugs are now [DOING] or [DONE]\n"
        "   ☐  Live URL works for everyone in the group on their own laptop",
        color=ACCENT_SOFT, size=15,
    )


def main():
    prs = new_prs()
    slides = [
        slide_title,                # 1 cover
        slide_workflow,             # 2
        slide_roles,                # 3
        slide_setup,                # 4
        slide_plan,                 # 5
        slide_naming,               # 6
        slide_glue,                 # 7
        slide_code_your_feature,    # 8
        slide_buglist,              # 9
        slide_streamlit_playground, # 10
        slide_pure_python,          # 11
        slide_requirements,         # 12
        slide_access,               # 13
        slide_passwords_3_8,        # 14
        slide_passwords_9_11,       # 15
        slide_passwords_12_16,      # 16
        slide_ship,                 # 17
        slide_iterate,              # 18
    ]
    total = len(slides) - 1
    slides[0](prs, total)
    for i, fn in enumerate(slides[1:], start=1):
        fn(prs, page=i, total=total)
    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT)
    print(f"Wrote: {OUTPUT}")


if __name__ == "__main__":
    main()
